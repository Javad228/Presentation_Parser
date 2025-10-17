#!/usr/bin/env python3
# -*- coding: utf-8 -*-

"""
pptx_spatial_map.py
- Extracts a spatial JSON map from a .pptx
- Classifies 'badge' text (single digit or "!") with a visible background as ICON
- Auto-renders per-slide preview PNGs with a legend

Usage:
  python pptx_spatial_map.py input.pptx
Outputs:
  input.spatial.json
  input_previews/slide_XX.png
"""

import json, os, re, sys, uuid, argparse
from datetime import datetime
try:
    from flask import Flask, request, redirect, url_for, send_from_directory, jsonify, render_template_string, abort
except Exception:
    # Flask is optional for CLI usage; import lazily in serve mode
    Flask = None  # type: ignore
    request = None  # type: ignore
    redirect = None  # type: ignore
    url_for = None  # type: ignore
    send_from_directory = None  # type: ignore
    jsonify = None  # type: ignore
    render_template_string = None  # type: ignore
    abort = None  # type: ignore
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

NS = {
    "a": "http://schemas.openxmlformats.org/drawingml/2006/main",
    "p": "http://schemas.openxmlformats.org/presentationml/2006/main",
    "c": "http://schemas.openxmlformats.org/drawingml/2006/chart",
    "dgm": "http://schemas.openxmlformats.org/drawingml/2006/diagram",
}

# ---------- small helpers ----------
def tag_name(shape):
    try:
        t = shape._element.tag
        return t.rsplit('}', 1)[-1] if '}' in t else t
    except Exception:
        return ""

def rel_bbox(left, top, width, height, total_w, total_h):
    def r(v, denom): return round(float(v)/float(denom), 6) if denom else 0.0
    return [r(left, total_w), r(top, total_h), r(width, total_w), r(height, total_h)]

def _safe_has_text_frame(shape):
    try:
        return bool(getattr(shape, "has_text_frame", False) and shape.text_frame)
    except Exception:
        return False

def _extract_text(shape, max_chars=2000):
    if not _safe_has_text_frame(shape):
        return ""
    try:
        return "".join([p.text or "" for p in shape.text_frame.paragraphs])[:max_chars]
    except Exception:
        return ""

def has_text(shape):
    try:
        return bool(_extract_text(shape).strip())
    except Exception:
        return False

def is_chart(shape):
    try:
        _ = shape.chart
        return True
    except Exception:
        pass
    try:
        gd = shape._element.xpath(".//a:graphicData", namespaces=NS)
        return any(el.get("uri") == NS["c"] for el in gd)
    except Exception:
        return False

def is_table(shape):
    try:
        _ = shape.table
        return True
    except Exception:
        return False

def is_smart_art(shape):
    try:
        gd = shape._element.xpath(".//a:graphicData", namespaces=NS)
        return any(el.get("uri") == NS["dgm"] for el in gd)
    except Exception:
        return False

def _has_visible_fill_or_line(shape):
    """Return True if the shape has a visible line or fill (not a:noFill)."""
    try:
        el = shape._element
        # line
        for l in el.xpath(".//a:ln", namespaces=NS):
            if l.find("./a:noFill", namespaces=NS) is None:
                return True
        # fill on spPr
        spPr = el.find(".//a:spPr", namespaces=NS)
        if spPr is not None and spPr.find("./a:noFill", namespaces=NS) is None:
            if any(spPr.find(tag, namespaces=NS) is not None
                   for tag in ("./a:solidFill", "./a:gradFill", "./a:pattFill")):
                return True
        return False
    except Exception:
        return False

def _extract_text_style(shape):
    """Return {'font_pt': float|None, 'bold': bool|None} aggregated from runs.
    Chooses the maximum font size seen; bold is True if any run is bold.
    """
    try:
        if not _safe_has_text_frame(shape):
            return {"font_pt": None, "bold": None}
        max_pt = None
        any_bold = None
        for p in shape.text_frame.paragraphs:
            for r in getattr(p, 'runs', []) or []:
                fnt = getattr(r, 'font', None)
                if fnt is None:
                    continue
                sz = getattr(fnt, 'size', None)
                if sz is not None:
                    try:
                        pt = float(getattr(sz, 'pt', None) or 0.0)
                        if pt > 0.0:
                            if (max_pt is None) or (pt > max_pt):
                                max_pt = pt
                    except Exception:
                        pass
                b = getattr(fnt, 'bold', None)
                if b is not None:
                    any_bold = bool(b) if any_bold is None else (any_bold or bool(b))
        return {"font_pt": max_pt, "bold": any_bold}
    except Exception:
        return {"font_pt": None, "bold": None}

# ---------- classification ----------
_BADGE_ICON_RE = re.compile(r"^(?:[0-9]{1,2}|!)$")
def _rel_area(shape, slide_hints):
    try:
        return (float(shape.width) * float(shape.height)) / (float(slide_hints["w"]) * float(slide_hints["h"]))
    except Exception:
        return 0.0


def classify(shape, slide_hints):
    t = tag_name(shape)

    # placeholders -> text
    try:
        if getattr(shape, "is_placeholder", False):
            return "text"
    except Exception:
        pass

    if t == "grpSp":
        return "group"
    if t == "cxnSp":
        return "connector"
    if t == "pic":
        return "image"

    if t == "graphicFrame":
        if is_chart(shape):     return "chart"
        if is_table(shape):     return "table"
        if is_smart_art(shape): return "figure"
        return "shape"

    if t == "sp":
        # Connector-like auto shapes (structural, not heuristic)
        try:
            connector_types = {
                MSO_SHAPE_TYPE.LINE, MSO_SHAPE_TYPE.STRAIGHT_CONNECTOR_1,
                MSO_SHAPE_TYPE.BENT_CONNECTOR_2, MSO_SHAPE_TYPE.BENT_CONNECTOR_3,
                MSO_SHAPE_TYPE.BENT_CONNECTOR_4, MSO_SHAPE_TYPE.BENT_CONNECTOR_5,
                MSO_SHAPE_TYPE.CURVED_CONNECTOR_2, MSO_SHAPE_TYPE.CURVED_CONNECTOR_3,
                MSO_SHAPE_TYPE.CURVED_CONNECTOR_4, MSO_SHAPE_TYPE.CURVED_CONNECTOR_5,
                MSO_SHAPE_TYPE.CURVE
            }
            if shape.auto_shape_type in connector_types:
                return "connector"
        except (AttributeError, KeyError):
            pass

        if has_text(shape):
            return "text"
        if _has_visible_fill_or_line(shape):
            return "shape"
        return "unknown"

    return "unknown"


# ---------- JSON model ----------
def base_component(shape, slide_idx, ctype, slide_hints, z=0, group_id=None, forced_id=None):
    left, top, width, height = int(shape.left), int(shape.top), int(shape.width), int(shape.height)
    rel = rel_bbox(left, top, width, height, slide_hints["w"], slide_hints["h"])
    cid = forced_id if forced_id else f"s{slide_idx}_c{shape.shape_id}"
    comp = {
        "id": cid,
        "type": ctype,
        "bbox_emus": [left, top, width, height],
        "bbox_rel": rel,
        "z": z,
        "group_id": group_id,
        "debug": {"tag": tag_name(shape)}
    }
    if ctype in ("text",):
        style = _extract_text_style(shape)
        comp["text_style"] = style
    return comp

def walk_group_children(group_shape, slide_idx, items, slide_hints, parent_id, z_start=0):
    z = z_start
    for shp in group_shape.shapes:
        ctype = classify(shp, slide_hints)
        if ctype == "group":
            gid = f"s{slide_idx}_g{shp.shape_id}"
            items.append(base_component(shp, slide_idx, "figure", slide_hints, z=z, group_id=parent_id, forced_id=gid))
            z = walk_group_children(shp, slide_idx, items, slide_hints, gid, z_start=z+1)
        else:
            items.append(base_component(shp, slide_idx, ctype, slide_hints, z=z, group_id=parent_id))
            z += 1
    return z

def walk_shapes(slide, slide_idx, items, slide_hints):
    z = 0
    for shp in slide.shapes:
        ctype = classify(shp, slide_hints)
        if ctype == "group":
            gid = f"s{slide_idx}_g{shp.shape_id}"
            items.append(base_component(shp, slide_idx, "figure", slide_hints, z=z, forced_id=gid))
            z = walk_group_children(shp, slide_idx, items, slide_hints, gid, z_start=z+1)
        else:
            items.append(base_component(shp, slide_idx, ctype, slide_hints, z=z))
            z += 1
    return items

def _iou(b1, b2):
    # bboxes in relative coords [x,y,w,h]
    x1,y1,w1,h1 = b1; x2,y2,w2,h2 = b2
    ax1,ay1,ax2,ay2 = x1, y1, x1+w1, y1+h1
    bx1,by1,bx2,by2 = x2, y2, x2+w2, y2+h2
    ix1, iy1 = max(ax1, bx1), max(ay1, by1)
    ix2, iy2 = min(ax2, bx2), min(ay2, by2)
    iw = max(0.0, ix2 - ix1)
    ih = max(0.0, iy2 - iy1)
    inter = iw * ih
    a = w1 * h1; b = w2 * h2
    union = a + b - inter if (a + b - inter) > 0 else 1e-9
    return inter / union

def _suppress_icon_backgrounds(items):
    """Relabel background rectangles that sit behind short-text icons."""
    # index shapes by group
    by_group = {}
    for i, c in enumerate(items):
        by_group.setdefault(c.get("group_id"), []).append(i)

    for grp, idxs in by_group.items():
        icons = [i for i in idxs if items[i]["type"] == "icon"]
        shapes = [i for i in idxs if items[i]["type"] == "shape"]
        if not icons or not shapes:
            continue
        for ii in icons:
            ibox = items[ii]["bbox_rel"]
            iarea = ibox[2] * ibox[3]
            for sj in shapes:
                sbox = items[sj]["bbox_rel"]
                sarea = sbox[2] * sbox[3]
                # Heuristic: strong overlap and comparable size (bg slightly smaller or larger)
                if _iou(ibox, sbox) >= 0.55 and 0.4 <= (sarea / iarea) <= 1.6:
                    items[sj]["type"] = "icon_bg"    # mark as background


def process_pptx(path):
    prs = Presentation(path)
    out = {"file": path, "slides": []}
    for idx, slide in enumerate(prs.slides):
        hints = {"w": int(prs.slide_width), "h": int(prs.slide_height)}
        items = []
        walk_shapes(slide, idx, items, hints)
        _suppress_icon_backgrounds(items)   # <<< NEW: hide badge backgrounds
        out["slides"].append({
            "index": idx,
            "canvas": {"w_emus": hints["w"], "h_emus": hints["h"]},
            "components": items
        })
    return out


# ---------- previews ----------
TYPE_COLORS = {
    "text": (0,0,0),
    "table": (0,0,0),
    "chart": (0,0,0),
    "image": (0,0,0),
    "shape": (0,0,0),
    "connector": (0,0,0),
    "figure": (0,0,0),
    "unknown": (0,0,0),
}

TYPE_CODE = {
    "text": 1,
    "image": 2,
    "table": 3,
    "chart": 4,
    "shape": 5,
    "connector": 6,
    "figure": 7,
    "unknown": 8,
}

def draw_previews(mapping, out_dir, img_width=1600, stroke=3, draw_labels=True, show_legend=True, show_groups=True):
    from PIL import Image, ImageDraw, ImageFont
    os.makedirs(out_dir, exist_ok=True)

    def want(comp):
        return True

    for slide in mapping["slides"]:
        W_emus, H_emus = slide["canvas"]["w_emus"], slide["canvas"]["h_emus"]
        aspect = H_emus / float(W_emus) if W_emus else 9/16.0
        W = int(img_width)
        H = max(1, int(round(W*aspect)))

        img = Image.new("RGB", (W, H), (255, 255, 255))
        draw = ImageDraw.Draw(img)
        try:
            font = ImageFont.load_default()
        except Exception:
            font = None

        for comp in slide["components"]:
            if not want(comp):
                continue
            t = comp.get("type", "unknown")
            color = (0, 0, 0)
            x, y, w, h = comp["bbox_rel"]
            x0, y0 = int(round(x * W)), int(round(y * H))
            x1, y1 = int(round((x + w) * W)), int(round((y + h) * H))
            for s in range(stroke):
                draw.rectangle([x0 - s, y0 - s, x1 + s, y1 + s], outline=color, width=1)
            if draw_labels:
                code = str(TYPE_CODE.get(t, 8))
                pad = 3
                tw = max(6 * len(code), 14)
                th = 12
                draw.rectangle([x0, y0, x0 + tw + 2 * pad, y0 + th + 2 * pad], fill=(255, 255, 255))
                draw.text((x0 + pad, y0 + pad), code, fill=color, font=font)

        if show_legend:
            pad = 8
            cx, cy = pad, pad
            for name in ["text","image","table","chart","shape","connector","figure","unknown"]:
                code = TYPE_CODE[name]
                draw.text((cx, cy), f"{code}  {name}", fill=(0,0,0), font=font)
                cy += 16

        out_path = os.path.join(out_dir, f"slide_{slide['index']:02d}.png")
        img.save(out_path, "PNG")


# ---------- job mapping helpers ----------

def _load_job_mapping(job_dir: str):
    json_files = [f for f in os.listdir(job_dir) if f.endswith('.spatial.json')]
    if not json_files:
        raise FileNotFoundError('No mapping JSON found')
    json_file = json_files[0]
    with open(os.path.join(job_dir, json_file), 'r', encoding='utf-8') as f:
        mapping = json.load(f)
    return mapping, json_file


def _save_job_mapping(job_dir: str, mapping: dict, json_file: str) -> None:
    path = os.path.join(job_dir, json_file)
    with open(path, 'w', encoding='utf-8') as f:
        json.dump(mapping, f, indent=2)


# ---------- web server (restored) ----------

def create_app():
    if Flask is None:
        raise RuntimeError("Flask is not installed. Please install dependencies and try again.")

    _ensure_web_dirs()

    app = Flask(__name__)
    app.config["MAX_CONTENT_LENGTH"] = 128 * 1024 * 1024

    INDEX_HTML = """
    <!doctype html>
    <html>
    <head>
      <meta charset=\"utf-8\"/>
      <title>PPTX Spatial Map</title>
      <style>
        body { font-family: system-ui, -apple-system, Segoe UI, Roboto, sans-serif; margin: 24px; }
        .card { border: 1px solid #e2e8f0; border-radius: 8px; padding: 16px; max-width: 760px; }
        .row { margin-top: 16px; }
        input[type=file] { padding: 8px; border: 1px solid #cbd5e1; border-radius: 6px; width: 100%; }
        button { background: #2563eb; color: white; border: none; border-radius: 6px; padding: 10px 14px; cursor: pointer; }
        .jobs { margin-top: 32px; max-width: 760px; }
        .job { padding: 10px 0; border-bottom: 1px solid #e2e8f0; }
      </style>
    </head>
    <body>
      <h2>PPTX Spatial Map</h2>
      <div class=\"card\">
        <form action=\"{{ url_for('upload') }}\" method=\"post\" enctype=\"multipart/form-data\">
          <label for=\"file\">Upload a .pptx file</label>
          <div class=\"row\"><input id=\"file\" name=\"file\" type=\"file\" accept=\".pptx\" required /></div>
          <div class=\"row\"><button type=\"submit\">Process</button></div>
        </form>
      </div>
      {% if jobs %}
      <div class=\"jobs\">
        <h3>Recent Jobs</h3>
        {% for j in jobs %}
          <div class=\"job\">
            <a href=\"{{ url_for('view_job', job_id=j['job_id']) }}\">{{ j['pptx_name'] }}</a>
            <span> · {{ j['num_slides'] }} slides</span>
          </div>
        {% endfor %}
      </div>
      {% endif %}
    </body>
    </html>
    """

    JOB_HTML = """
    <!doctype html>
    <html>
    <head>
      <meta charset=\"utf-8\"/>
      <title>Job {{ job_id }} - PPTX Spatial Map</title>
      <style>
        body { font-family: system-ui, -apple-system, Segoe UI, Roboto, sans-serif; margin: 24px; }
        .row { margin: 16px 0; }
        .grid { display: grid; grid-template-columns: repeat(auto-fill, minmax(360px, 1fr)); gap: 16px; }
        .card { border: 1px solid #e2e8f0; border-radius: 8px; padding: 12px; }
        img { width: 100%; height: auto; border-radius: 6px; border: 1px solid #e2e8f0; }
        a.button { background: #2563eb; color: white; padding: 8px 12px; border-radius: 6px; text-decoration: none; }
      </style>
    </head>
    <body>
      <a href=\"{{ url_for('index') }}\">← Back</a>
      <h2>{{ pptx_name }}</h2>
      <div class=\"row\">
        <a class=\"button\" href=\"{{ url_for('download_json', job_id=job_id) }}\">Download JSON</a>
      </div>
      <div class=\"grid\">
        {% for i in range(num_slides) %}
          <div class=\"card\">
            <div style=\"display:flex; justify-content: space-between; align-items:center; margin-bottom:6px;\">
              <strong>Slide {{ i }}</strong>
              <a class=\"button\" href=\"{{ url_for('edit_slide', job_id=job_id, slide_index=i) }}\">Edit</a>
            </div>
            <img loading=\"lazy\" src=\"{{ url_for('preview_image', job_id=job_id, filename='slide_' + ('%02d' % i) + '.png') }}\" alt=\"slide {{ i }}\" />
          </div>
        {% endfor %}
      </div>
    </body>
    </html>
    """

    EDITOR_HTML = """
    <!doctype html>
    <html>
    <head>
      <meta charset=\"utf-8\"/>
      <title>Edit Slide {{ slide_index }} - {{ pptx_name }}</title>
      <style>
        body { font-family: system-ui, -apple-system, Segoe UI, Roboto, sans-serif; margin: 16px; }
        .toolbar { display:flex; gap:8px; align-items:center; margin-bottom: 12px; }
        .button { background:#2563eb; color:#fff; border:none; border-radius:6px; padding:8px 12px; text-decoration:none; cursor:pointer; }
        .layout { display:grid; grid-template-columns: 1fr 1fr; gap:16px; align-items:start; }
        .pane { border:1px solid #e2e8f0; border-radius:8px; padding:8px; }
        .stage { position:relative; width:100%; background:#fff; aspect-ratio: {{ canvas_w }} / {{ canvas_h }}; overflow:hidden; }
        .rect { position:absolute; border:2px solid #111827; background: rgba(17,24,39,0.05); box-sizing: border-box; }
        .rect[data-type=text] { border-color:#1e90ff; background: rgba(30,144,255,0.08); }
        .rect[data-type=icon] { border-color:#9933ff; background: rgba(153,51,255,0.08); }
        .rect[data-type=image] { border-color:#ff8c00; background: rgba(255,140,0,0.10); }
        .rect[data-type=shape] { border-color:#008b8b; background: rgba(0,139,139,0.08); }
        .rect[data-type=table] { border-color:#228b22; background: rgba(34,139,34,0.10); }
        .rect[data-type=chart] { border-color:#dc143c; background: rgba(220,20,60,0.10); }
        .rect[data-type=connector] { border-color:#787878; background: rgba(120,120,120,0.10); }
        .rect[data-type=figure] { border-color:#bdb76b; background: rgba(189,183,107,0.10); }
        .rect[data-type=unknown] { border-color:#666666; background: rgba(100,100,100,0.10); }
        .rect.selected { outline: 2px dashed #111827; }
        .handle { position:absolute; width:12px; height:12px; background:#111827; border-radius:2px; }
        .handle.tl { left:-6px; top:-6px; cursor:nwse-resize; }
        .handle.tr { right:-6px; top:-6px; cursor:nesw-resize; }
        .handle.bl { left:-6px; bottom:-6px; cursor:nesw-resize; }
        .handle.br { right:-6px; bottom:-6px; cursor:nwse-resize; }
        .meta { font-size:12px; color:#334155; margin-top:6px; }
        .legend { display:flex; flex-wrap:wrap; gap:10px; margin-top:8px; }
        .legend .swatch { display:inline-block; width:14px; height:14px; border:2px solid currentColor; vertical-align:middle; margin-right:6px; }
        .slideimg { width:100%; height:auto; border-radius:6px; border:1px solid #e2e8f0; }
        .tooltip { position:absolute; background:#111827; color:#fff; padding:4px 6px; border-radius:4px; font-size:12px; pointer-events:none; transform: translate(-50%, -120%); display:none; white-space:nowrap; }
      </style>
    </head>
    <body>
      <div class=\"toolbar\">
        <a class=\"button\" href=\"{{ url_for('view_job', job_id=job_id) }}\">← Back</a>
        <a class=\"button\" href=\"{{ url_for('download_json', job_id=job_id) }}\">Download JSON</a>
        <span style=\"margin-left:auto\"></span>
        <a class=\"button\" href=\"{{ url_for('edit_slide', job_id=job_id, slide_index=prev_index) }}\">Prev</a>
        <a class=\"button\" href=\"{{ url_for('edit_slide', job_id=job_id, slide_index=next_index) }}\">Next</a>
        <button id=\"deleteBtn\" class=\"button\" style=\"background:#ef4444\">Delete selected</button>
      </div>

      <div class=\"layout\">
        <div class=\"pane\">
          <div id=\"stage\" class=\"stage\"></div>
          <div id=\"meta\" class=\"meta\"></div>
          <div class=\"legend\">
            <span>1 text</span>
            <span>2 image</span>
            <span>3 table</span>
            <span>4 chart</span>
            <span>5 shape</span>
            <span>6 connector</span>
            <span>7 figure</span>
            <span>8 unknown</span>
          </div>
        </div>
        <div class=\"pane\">
          <img class=\"slideimg\" src=\"{{ url_for('slide_image', job_id=job_id, slide_index=slide_index) }}\" alt=\"Slide image\"/>
        </div>
      </div>

      <script>
      const jobId = {{ job_id|tojson }};
      const slideIndex = {{ slide_index|tojson }};
      let slideData = null;
      let selectedId = null;

      function px(n) { return n + 'px'; }

      async function loadSlide() {
        const res = await fetch(`/api/job/${jobId}/slide/${slideIndex}`);
        if (!res.ok) { alert('Failed to load slide'); return; }
        slideData = await res.json();
        renderStage();
      }

      function createHandle(cls) {
        const h = document.createElement('div');
        h.className = 'handle ' + cls;
        return h;
      }

      function renderStage() {
        const stage = document.getElementById('stage');
        stage.innerHTML = '';
        const tip = document.createElement('div');
        tip.id = 'tooltip'; tip.className = 'tooltip';
        stage.appendChild(tip);
        const W = stage.clientWidth;
        const H = stage.clientHeight;
        const comps = slideData.components;
        document.getElementById('meta').textContent = `${comps.length} components`;
        for (const c of comps) {
          if (!c.bbox_rel) continue;
          const [x,y,w,h] = c.bbox_rel;
          const el = document.createElement('div');
          el.className = 'rect';
          el.dataset.id = c.id;
          el.dataset.type = c.type || 'unknown';
          el.style.left = Math.round(x * W) + 'px';
          el.style.top = Math.round(y * H) + 'px';
          el.style.width = Math.max(2, Math.round(w * W)) + 'px';
          el.style.height = Math.max(2, Math.round(h * H)) + 'px';
          const tl = createHandle('tl'); const tr = createHandle('tr');
          const bl = createHandle('bl'); const br = createHandle('br');
          el.appendChild(tl); el.appendChild(tr); el.appendChild(bl); el.appendChild(br);
          // Numeric badge (code) in top-left for accessibility
          const code = (c.type && ({text:1,image:2,table:3,chart:4,shape:5,connector:6,figure:7,unknown:8})[c.type]) || 8;
          const badge = document.createElement('div');
          badge.style.position = 'absolute';
          badge.style.left = '0px'; badge.style.top = '0px';
          badge.style.background = 'rgba(255,255,255,0.9)';
          badge.style.border = '1px solid #111827';
          badge.style.borderBottomRightRadius = '6px';
          badge.style.padding = '1px 4px';
          badge.style.fontSize = '12px';
          badge.textContent = String(code);
          el.appendChild(badge);
          // Hover tooltip for text/icon font style
          el.addEventListener('mousemove', (e) => {
            const st = c.text_style || {};
            const lines = [];
            if (st.font_pt) lines.push(`font ${st.font_pt.toFixed(1)}pt`);
            if (typeof st.bold === 'boolean') lines.push(st.bold ? 'bold' : 'normal');
            tip.textContent = lines.join(' · ') || '';
            tip.style.left = e.clientX - stage.getBoundingClientRect().left + 'px';
            tip.style.top = e.clientY - stage.getBoundingClientRect().top + 'px';
            tip.style.display = tip.textContent ? 'block' : 'none';
          });
          el.addEventListener('mouseleave', () => { tip.style.display = 'none'; });
          stage.appendChild(el);
        }
        bindInteractions();
      }

      function bindInteractions() {
        const stage = document.getElementById('stage');
        const W = stage.clientWidth; const H = stage.clientHeight;
        let drag = null;
        stage.querySelectorAll('.rect').forEach(el => {
          el.addEventListener('mousedown', (e) => {
            document.querySelectorAll('.rect.selected').forEach(r => r.classList.remove('selected'));
            el.classList.add('selected');
            const rect = el.getBoundingClientRect();
            const start = { x: e.clientX, y: e.clientY, left: rect.left, top: rect.top, width: rect.width, height: rect.height };
            const isHandle = e.target.classList.contains('handle');
            const handle = isHandle ? (e.target.classList.contains('tl') ? 'tl' : e.target.classList.contains('tr') ? 'tr' : e.target.classList.contains('bl') ? 'bl' : 'br') : null;
            drag = { target: el, start, handle };
            window.addEventListener('mousemove', onMove);
            window.addEventListener('mouseup', onUp, { once: true });
            e.preventDefault();
          });
        });
        stage.addEventListener('mousedown', (e) => { if (!e.target.closest('.rect')) document.querySelectorAll('.rect.selected').forEach(r => r.classList.remove('selected')); });

        function onMove(e) {
          if (!drag) return;
          const dx = e.clientX - drag.start.x;
          const dy = e.clientY - drag.start.y;
          let left = drag.start.left + dx - stage.getBoundingClientRect().left;
          let top = drag.start.top + dy - stage.getBoundingClientRect().top;
          let width = drag.start.width;
          let height = drag.start.height;
          if (!drag.handle) {
            left = Math.max(0, Math.min(W - width, left));
            top = Math.max(0, Math.min(H - height, top));
          } else {
            if (drag.handle.includes('t')) { height = Math.max(4, Math.min(H, drag.start.height - dy)); top = Math.max(0, Math.min(drag.start.top - stage.getBoundingClientRect().top + drag.start.height - 4, drag.start.top + dy - stage.getBoundingClientRect().top)); }
            if (drag.handle.includes('b')) { height = Math.max(4, Math.min(H, drag.start.height + dy)); }
            if (drag.handle.includes('l')) { width = Math.max(4, Math.min(W, drag.start.width - dx)); left = Math.max(0, Math.min(drag.start.left - stage.getBoundingClientRect().left + drag.start.width - 4, drag.start.left + dx - stage.getBoundingClientRect().left)); }
            if (drag.handle.includes('r')) { width = Math.max(4, Math.min(W, drag.start.width + dx)); }
          }
          drag.target.style.left = left + 'px';
          drag.target.style.top = top + 'px';
          drag.target.style.width = width + 'px';
          drag.target.style.height = height + 'px';
        }
        async function onUp() {
          if (!drag) return;
          window.removeEventListener('mousemove', onMove);
          const el = drag.target; drag = null;
          const left = parseInt(el.style.left)/W; const top = parseInt(el.style.top)/H;
          const width = parseInt(el.style.width)/W; const height = parseInt(el.style.height)/H;
          await fetch(`/api/job/${jobId}/slide/${slideIndex}/component/${encodeURIComponent(el.dataset.id)}`, {
            method: 'POST', headers: { 'Content-Type': 'application/json' }, body: JSON.stringify({ bbox_rel: [left, top, width, height] })
          });
        }
        document.getElementById('deleteBtn').onclick = async () => {
          const sel = document.querySelector('.rect.selected');
          if (!sel) return;
          if (!confirm('Delete selected component?')) return;
          const ok = await fetch(`/api/job/${jobId}/slide/${slideIndex}/component/${encodeURIComponent(sel.dataset.id)}`, { method: 'DELETE' });
          if (ok) loadSlide();
        };
      }

      window.addEventListener('resize', () => renderStage());
      loadSlide();
      </script>
    </body>
    </html>
    """

    def _list_jobs(limit=20):
        jobs = []
        if not os.path.isdir(JOBS_DIR):
            return jobs
        for job_id in sorted(os.listdir(JOBS_DIR), reverse=True):
            job_dir = os.path.join(JOBS_DIR, job_id)
            if not os.path.isdir(job_dir):
                continue
            json_files = [f for f in os.listdir(job_dir) if f.endswith('.spatial.json')]
            if not json_files:
                continue
            json_file = json_files[0]
            try:
                with open(os.path.join(job_dir, json_file), 'r', encoding='utf-8') as f:
                    data = json.load(f)
                pptx_name = os.path.basename(data.get('file', 'presentation.pptx'))
                num_slides = len((data or {}).get('slides', []))
                jobs.append({"job_id": job_id, "pptx_name": pptx_name, "num_slides": num_slides})
            except Exception:
                continue
            if len(jobs) >= limit:
                break
        return jobs

    @app.get("/")
    def index():
        return render_template_string(INDEX_HTML, jobs=_list_jobs())

    @app.post("/upload")
    def upload():
        if "file" not in request.files:
            abort(400, "No file part in request")
        file = request.files["file"]
        if file.filename == "":
            abort(400, "No file selected")
        if not _allowed_file(file.filename):
            abort(400, "Only .pptx files are supported")
        temp_path = os.path.join(UPLOAD_DIR, uuid.uuid4().hex + ".pptx")
        file.save(temp_path)
        job = _process_upload_to_job(temp_path, file.filename)
        return redirect(url_for("view_job", job_id=job["job_id"]))

    @app.get("/job/<job_id>")
    def view_job(job_id):
        job_dir = os.path.join(JOBS_DIR, job_id)
        if not os.path.isdir(job_dir):
            abort(404)
        json_files = [f for f in os.listdir(job_dir) if f.endswith('.spatial.json')]
        if not json_files:
            abort(404)
        json_file = json_files[0]
        with open(os.path.join(job_dir, json_file), 'r', encoding='utf-8') as f:
            data = json.load(f)
        pptx_name = os.path.basename(data.get('file', 'presentation.pptx'))
        num_slides = len((data or {}).get('slides', []))
        return render_template_string(JOB_HTML, job_id=job_id, pptx_name=pptx_name, num_slides=num_slides)

    @app.get("/job/<job_id>/edit/<int:slide_index>")
    def edit_slide(job_id, slide_index):
        job_dir = os.path.join(JOBS_DIR, job_id)
        if not os.path.isdir(job_dir):
            abort(404)
        mapping, _ = _load_job_mapping(job_dir)
        slides = mapping.get('slides', [])
        if slide_index < 0 or slide_index >= len(slides):
            abort(404)
        canvas = slides[slide_index].get('canvas', {})
        pptx_name = os.path.basename(mapping.get('file', 'presentation.pptx'))
        return render_template_string(EDITOR_HTML,
                                      job_id=job_id,
                                      slide_index=slide_index,
                                      num_slides=len(slides),
                                      prev_index=0 if slide_index - 1 < 0 else slide_index - 1,
                                      next_index=(len(slides) - 1) if slide_index + 1 >= len(slides) else slide_index + 1,
                                      pptx_name=pptx_name,
                                      canvas_w=canvas.get('w_emus', 1600),
                                      canvas_h=canvas.get('h_emus', 900))

    @app.get("/api/job/<job_id>/slide/<int:slide_index>")
    def get_slide(job_id, slide_index):
        job_dir = os.path.join(JOBS_DIR, job_id)
        mapping, _ = _load_job_mapping(job_dir)
        slides = mapping.get('slides', [])
        if slide_index < 0 or slide_index >= len(slides):
            abort(404)
        return jsonify(slides[slide_index])

    @app.post("/api/job/<job_id>/slide/<int:slide_index>/component/<cid>")
    def update_component(job_id, slide_index, cid):
        job_dir = os.path.join(JOBS_DIR, job_id)
        mapping, json_file = _load_job_mapping(job_dir)
        slides = mapping.get('slides', [])
        if slide_index < 0 or slide_index >= len(slides):
            abort(404)
        payload = request.get_json(force=True, silent=True) or {}
        bbox = payload.get('bbox_rel')
        if not (isinstance(bbox, list) and len(bbox) == 4 and all(isinstance(v, (int, float)) for v in bbox)):
            abort(400, 'bbox_rel must be [x,y,w,h]')
        x,y,w,h = bbox
        x = max(0.0, min(1.0, float(x)))
        y = max(0.0, min(1.0, float(y)))
        w = max(0.0, min(1.0 - x, float(w)))
        h = max(0.0, min(1.0 - y, float(h)))
        updated = None
        for comp in slides[slide_index].get('components', []):
            if comp.get('id') == cid:
                comp['bbox_rel'] = [x,y,w,h]
                updated = comp
                break
        if updated is None:
            abort(404)
        _save_job_mapping(job_dir, mapping, json_file)
        return jsonify(updated)

    @app.delete("/api/job/<job_id>/slide/<int:slide_index>/component/<cid>")
    def delete_component(job_id, slide_index, cid):
        job_dir = os.path.join(JOBS_DIR, job_id)
        mapping, json_file = _load_job_mapping(job_dir)
        slides = mapping.get('slides', [])
        if slide_index < 0 or slide_index >= len(slides):
            abort(404)
        comps = slides[slide_index].get('components', [])
        before = len(comps)
        comps[:] = [c for c in comps if c.get('id') != cid]
        if len(comps) == before:
            abort(404)
        _save_job_mapping(job_dir, mapping, json_file)
        return ('', 204)

    @app.get("/job/<job_id>/slide_image/<int:slide_index>")
    def slide_image(job_id, slide_index):
        job_dir = os.path.join(JOBS_DIR, job_id)
        img_path = _find_real_slide_image(job_dir, slide_index)
        if not img_path:
            try:
                pptx_files = [f for f in os.listdir(job_dir) if f.lower().endswith('.pptx')]
                if pptx_files:
                    _export_real_slide_images(os.path.join(job_dir, pptx_files[0]), job_dir)
                    img_path = _find_real_slide_image(job_dir, slide_index)
            except Exception:
                img_path = None
        if img_path:
            return send_from_directory(os.path.dirname(img_path), os.path.basename(img_path))
        filename = f"slide_{slide_index:02d}.png"
        previews_dirs = [d for d in os.listdir(job_dir) if d.endswith("_previews") and os.path.isdir(os.path.join(job_dir, d))]
        if not previews_dirs:
            abort(404)
        previews_dir = os.path.join(job_dir, previews_dirs[0])
        return send_from_directory(previews_dir, filename)

    @app.get("/job/<job_id>/json")
    def download_json(job_id):
        job_dir = os.path.join(JOBS_DIR, job_id)
        if not os.path.isdir(job_dir):
            abort(404)
        json_files = [f for f in os.listdir(job_dir) if f.endswith('.spatial.json')]
        if not json_files:
            abort(404)
        json_file = json_files[0]
        return send_from_directory(job_dir, json_file, as_attachment=True, download_name=json_file)

    @app.get("/job/<job_id>/previews/<path:filename>")
    def preview_image(job_id, filename):
        job_dir = os.path.join(JOBS_DIR, job_id)
        if not os.path.isdir(job_dir):
            abort(404)
        previews_dirs = [d for d in os.listdir(job_dir) if d.endswith("_previews") and os.path.isdir(os.path.join(job_dir, d))]
        if not previews_dirs:
            abort(404)
        previews_dir = os.path.join(job_dir, previews_dirs[0])
        return send_from_directory(previews_dir, filename)

    return app


def serve(host: str = "127.0.0.1", port: int = 5000) -> None:
    app = create_app()
    app.run(host=host, port=port, debug=False)


# ---------- web server ----------
WEB_ROOT = os.path.join(os.path.dirname(__file__), "webdata")
UPLOAD_DIR = os.path.join(WEB_ROOT, "uploads")
JOBS_DIR = os.path.join(WEB_ROOT, "jobs")


def _ensure_web_dirs():
    os.makedirs(UPLOAD_DIR, exist_ok=True)
    os.makedirs(JOBS_DIR, exist_ok=True)


def _allowed_file(filename):
    return "." in filename and filename.lower().rsplit(".", 1)[-1] in {"pptx"}


def _new_job_id():
    return datetime.utcnow().strftime("%Y%m%d-%H%M%S-") + uuid.uuid4().hex[:8]


def _process_upload_to_job(src_path, original_name):
    """Create a job folder, process PPTX, render previews, write JSON; return job info."""
    job_id = _new_job_id()
    job_dir = os.path.join(JOBS_DIR, job_id)
    os.makedirs(job_dir, exist_ok=True)

    # Copy or move upload into job
    base_name = os.path.basename(original_name)
    stored_pptx = os.path.join(job_dir, base_name)
    try:
        # Prefer atomic rename within same filesystem when possible
        os.replace(src_path, stored_pptx)
    except Exception:
        # Fallback to copy if replace fails
        import shutil
        shutil.copy2(src_path, stored_pptx)

    mapping = process_pptx(stored_pptx)

    # Write JSON
    json_path = os.path.join(job_dir, os.path.splitext(base_name)[0] + ".spatial.json")
    with open(json_path, "w", encoding="utf-8") as f:
        json.dump(mapping, f, indent=2)

    # Render previews inside job dir
    previews_dir = os.path.join(job_dir, os.path.splitext(base_name)[0] + "_previews")
    draw_previews(mapping, previews_dir, img_width=1200, stroke=2, draw_labels=True, show_legend=True, show_groups=True)

    # Try to export real slide images via LibreOffice, if available
    try:
        _export_real_slide_images(stored_pptx, job_dir)
    except Exception:
        pass

    return {
        "job_id": job_id,
        "job_dir": job_dir,
        "pptx_name": base_name,
        "json_file": os.path.basename(json_path),
        "previews_rel": os.path.basename(previews_dir),
        "num_slides": len(mapping.get("slides", [])),
    }


def _export_real_slide_images(pptx_path, job_dir):
    """Export real slide images.
    Strategy:
      1) Try LibreOffice to PNG for all slides
      2) If the count of PNGs < number of slides, fall back to PDF + pdftoppm
      3) Normalize filenames to slide-<1-based>.png for deterministic lookup
    Creates/uses job_dir/slides_png
    """
    import shutil, subprocess, tempfile
    out_dir = os.path.join(job_dir, "slides_png")
    os.makedirs(out_dir, exist_ok=True)

    # Determine expected slide count from PPTX
    try:
        expected = len(Presentation(pptx_path).slides)
    except Exception:
        expected = None

    def list_pngs():
        try:
            return sorted([f for f in os.listdir(out_dir) if f.lower().endswith('.png')])
        except Exception:
            return []

    # If we already have per-slide PNGs and they match expectation, still normalize names
    existing = list_pngs()

    soffice = shutil.which("soffice") or shutil.which("libreoffice")
    if not existing and soffice:
        try:
            subprocess.run([soffice, "--headless", "--convert-to", "png", "--outdir", out_dir, pptx_path],
                           check=True, stdout=subprocess.PIPE, stderr=subprocess.PIPE, timeout=180)
        except Exception:
            pass
        existing = list_pngs()

    # Fallback: convert to PDF then render pages via pdftoppm if needed
    pdftoppm = shutil.which("pdftoppm")
    if (not existing) or (expected and len(existing) < expected):
        if soffice and pdftoppm:
            try:
                with tempfile.TemporaryDirectory() as td:
                    pdf_path = os.path.join(td, "slides.pdf")
                    # Export PPTX to PDF
                    subprocess.run([soffice, "--headless", "--convert-to", "pdf", "--outdir", td, pptx_path],
                                   check=True, stdout=subprocess.PIPE, stderr=subprocess.PIPE, timeout=180)
                    # Sometimes LibreOffice preserves filename; find the pdf
                    if not os.path.isfile(pdf_path):
                        pdfs = [f for f in os.listdir(td) if f.lower().endswith('.pdf')]
                        if not pdfs:
                            raise RuntimeError("PDF not produced")
                        pdf_path = os.path.join(td, pdfs[0])
                    # Render PDF pages to PNGs at reasonable DPI
                    prefix = os.path.join(out_dir, "slide")
                    subprocess.run([pdftoppm, "-png", "-r", "144", pdf_path, prefix],
                                   check=True, stdout=subprocess.PIPE, stderr=subprocess.PIPE, timeout=180)
            except Exception:
                pass
        existing = list_pngs()

    # Normalize filenames deterministically to slide-<1-based>.png
    _normalize_slide_pngs(out_dir, expected)


def _normalize_slide_pngs(out_dir: str, expected: int | None) -> None:
    try:
        files = sorted([f for f in os.listdir(out_dir) if f.lower().endswith('.png')])
    except Exception:
        return
    if not files:
        return
    import re, os
    # Map files to numeric slide index if found in filename
    file_to_num: dict[str, int] = {}
    for f in files:
        nums = re.findall(r"(\d+)", f)
        if nums:
            try:
                n = int(nums[-1])  # use last number
                file_to_num[f] = n  # assume 1-based
            except Exception:
                pass
    used = set()
    # First, move files with explicit numbers to normalized names
    for f, n in file_to_num.items():
        target = os.path.join(out_dir, f"slide-{n:02d}.png")
        src = os.path.join(out_dir, f)
        if os.path.abspath(src) == os.path.abspath(target):
            used.add(f)
            continue
        # Avoid overwriting: remove existing target
        if os.path.exists(target):
            try:
                os.remove(target)
            except Exception:
                pass
        try:
            os.replace(src, target)
            used.add(f)
        except Exception:
            pass
    # Refresh list after moves
    try:
        remaining = sorted([f for f in os.listdir(out_dir) if f.lower().endswith('.png') and not re.fullmatch(r"slide-\d{2}\.png", f)])
    except Exception:
        remaining = []
    # Assign remaining files by sorted order
    # Determine max count to normalize to
    count = expected if (isinstance(expected, int) and expected > 0) else len(remaining)
    i = 1
    for f in remaining:
        # Find next unused normalized name
        while True:
            target = os.path.join(out_dir, f"slide-{i:02d}.png")
            i += 1
            if not os.path.exists(target):
                break
        src = os.path.join(out_dir, f)
        try:
            os.replace(src, target)
        except Exception:
            pass


def _find_real_slide_image(job_dir, slide_index):
    dirp = os.path.join(job_dir, 'slides_png')
    if not os.path.isdir(dirp):
        return None
    # Support both unpadded and zero-padded filenames
    unp = os.path.join(dirp, f"slide-{slide_index+1}.png")
    if os.path.isfile(unp):
        return unp
    pad = os.path.join(dirp, f"slide-{slide_index+1:02d}.png")
    if os.path.isfile(pad):
        return pad
    return None


def _crop_save_png(src_img_path, rel_bbox, out_path, pad_px=2):
    from PIL import Image
    with Image.open(src_img_path) as im:
        W, H = im.size
        x, y, w, h = rel_bbox
        x0 = max(0, int(round(x * W)) - pad_px)
        y0 = max(0, int(round(y * H)) - pad_px)
        x1 = min(W, int(round((x + w) * W)) + pad_px)
        y1 = min(H, int(round((y + h) * H)) + pad_px)
        if x1 <= x0 or y1 <= y0:
            return False
        crop = im.crop((x0, y0, x1, y1))
        os.makedirs(os.path.dirname(out_path), exist_ok=True)
        crop.save(out_path, 'PNG')
        return True

# ---------- main ----------
def main():
    parser = argparse.ArgumentParser(description="Extract spatial JSON and previews from a PPTX, or run a web server.")
    subparsers = parser.add_subparsers(dest="command")

    # Serve subcommand
    p_serve = subparsers.add_parser("serve", help="Run local web server")
    p_serve.add_argument("--host", default="127.0.0.1", help="Host to bind (default: 127.0.0.1)")
    p_serve.add_argument("--port", type=int, default=5000, help="Port to bind (default: 5000)")

    # CLI default: process a single file
    p_cli = subparsers.add_parser("process", help="Process a PPTX via CLI")
    p_cli.add_argument("input_pptx", help="Path to input .pptx")

    # Backward compatibility: allow calling without subcommand
    parser.add_argument("legacy_input", nargs="?", help=argparse.SUPPRESS)

    args = parser.parse_args()

    if args.command == "serve":
        serve(host=args.host, port=args.port)
        return

    in_path = None
    if args.command == "process":
        in_path = args.input_pptx
    elif getattr(args, "legacy_input", None):
        in_path = args.legacy_input
    else:
        parser.print_help()
        sys.exit(1)

    if not os.path.isfile(in_path):
        print(f"File not found: {in_path}")
        sys.exit(1)

    mapping = process_pptx(in_path)

    json_path = os.path.splitext(in_path)[0] + ".spatial.json"
    with open(json_path, "w", encoding="utf-8") as f:
        json.dump(mapping, f, indent=2)
    print(f"Wrote JSON: {json_path}")

    preview_dir = os.path.splitext(in_path)[0] + "_previews"
    draw_previews(mapping, preview_dir, img_width=1600, stroke=3, draw_labels=True, show_legend=True, show_groups=True)
    print(f"Wrote previews to: {preview_dir}")

if __name__ == "__main__":
    main()
